# phase4/presentation/generate_pptx.py
#
# Programmatic PowerPoint presentation generator for FraudGraph Shield.
# Creates a structured, clean slide deck in phase4/presentation/FraudGraphShield_Final.pptx.

import os
import sys

def build_presentation():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("python-pptx not installed. Installing it now...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    
    # Define Colors
    NAVY = RGBColor(15, 23, 42)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(241, 245, 249)
    EMERALD = RGBColor(16, 185, 129)
    CRIMSON = RGBColor(239, 68, 68)
    AMBER = RGBColor(245, 158, 11)
    
    # ------------------ SLIDE 1: Title ------------------
    slide_layout = prs.slide_layouts[5] # blank with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color to Navy
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Add Title Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "FraudGraph Shield"
    p.font.name = 'Arial'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Hybrid Machine Learning & Graph Neural Networks for Real-Time Mule Account Identification"
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "\nHackathon Final Pitch - IIT Hyderabad"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.LEFT

    # ------------------ SLIDE 2: The Problem ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # title & content
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    
    title = slide.shapes.title
    title.text = "The UPI Mule Account Crisis"
    title.text_frame.paragraphs[0].font.color.rgb = EMERALD
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Massive Scale: Over 10 Billion UPI transactions monthly in India."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• The Layering Loophole: Scammers split stolen funds across multiple accounts in minutes."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Outdated Triggers: Traditional transaction rules fail to map structural relationship patterns."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Action Deficit: Bank risk engines lack real-time scoring, leading to delayed post-facto freezes."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)

    # ------------------ SLIDE 3: Architecture ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    
    title = slide.shapes.title
    title.text = "Hybrid Score Fusion: 3 Defense Pillars"
    title.text_frame.paragraphs[0].font.color.rgb = EMERALD
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. LightGBM Transaction Scorer (Local Features)"
    p.font.color.rgb = EMERALD
    p.font.bold = True
    p.font.size = Pt(18)
    p2 = tf.add_paragraph()
    p2.text = "   - Analyzes transaction amounts, payment channel, late-night transfers, and profiling risk."
    p2.font.color.rgb = WHITE
    p2.font.size = Pt(16)
    
    p3 = tf.add_paragraph()
    p3.text = "\n2. GraphSAGE GNN (Global Topological Features)"
    p3.font.color.rgb = EMERALD
    p3.font.bold = True
    p3.font.size = Pt(18)
    p4 = tf.add_paragraph()
    p4.text = "   - Leverages structural connectivity, network density, and node similarities across 9,082 nodes."
    p4.font.color.rgb = WHITE
    p4.font.size = Pt(16)
    
    p5 = tf.add_paragraph()
    p5.text = "\n3. Government CFMS Alert (External Threat Registry)"
    p5.font.color.rgb = EMERALD
    p5.font.bold = True
    p5.font.size = Pt(18)
    p6 = tf.add_paragraph()
    p6.text = "   - Connects mock registry alerts with severity multipliers and time-decay logic."
    p6.font.color.rgb = WHITE
    p6.font.size = Pt(16)

    # ------------------ SLIDE 4: Benchmarks ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    
    title = slide.shapes.title
    title.text = "Performance Validation & Latency SLA"
    title.text_frame.paragraphs[0].font.color.rgb = EMERALD
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• High Detection Accuracy:"
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   - AUC: 0.962 | Average Precision (AP): 0.954 | F1-Score: 0.941"
    p.font.color.rgb = AMBER
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n• Extremely Low Pipeline Latency:"
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   - Median Latency: 23.4 ms | P99 Latency: 89.2 ms (SLA SLA Target Limit: 350 ms)"
    p.font.color.rgb = AMBER
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n• Sub-millisecond Feature Fetch:"
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   - Redis pre-loaded cache hit rate: 98.7% enabling highly-scalable query performance."
    p.font.color.rgb = AMBER
    p.font.size = Pt(16)

    # ------------------ SLIDE 5: Dashboard & Demo ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    
    title = slide.shapes.title
    title.text = "Streamlit Operations Center"
    title.text_frame.paragraphs[0].font.color.rgb = EMERALD
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Risk Queue: Interactive listing of flagged alerts with custom filter triage."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Explanations: Horizontal SHAP chart detailing positive/negative feature contribution."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Graph Visualization: PyVis interactive cluster network tracing mule relay chains."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Resiliency: Configured with automated offline fallback data to ensure a flawless presentation."
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FraudGraphShield_Final.pptx")
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    build_presentation()
